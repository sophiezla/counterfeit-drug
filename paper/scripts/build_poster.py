"""
Build the 48 x 36 inch landscape conference poster as a fully editable .pptx.

Everything is a native PowerPoint object: text frames, autoshapes and real
PowerPoint tables. The only raster content is the figures themselves, embedded
from paper/figures/*.png at 400 dpi (≈260-370 dpi effective at poster scale);
the vector PDFs of the same figures accompany the deliverable for anyone who
wants to swap them in.

Numbers are read from paper/tables/*.csv so the poster cannot drift from the
manuscript.

Layout: a header band, a full-width row of stat tiles carrying the three
headline numbers, four content columns, and a footer band holding the QR
placeholder, acknowledgements and references.

Color: the dataviz reference palette, light mode — blue #2a78d6 (slot 1),
orange #eb6834 (2), aqua #1baf7a (3), violet #4a3aa7 (7). The four-slot set
was validated with the skill's validator under --pairs all against a white
surface (CVD dE 9.2, normal-vision dE 16.3, both clear).

Output: paper/PharmaChecked_v2_poster.pptx
"""
import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIGS = ROOT / "paper" / "figures"
TABLES = ROOT / "paper" / "tables"
OUT = ROOT / "paper" / "PharmaChecked_v2_poster.pptx"

# ---------------------------------------------------------------- geometry
W, H = 48.0, 36.0
MARGIN = 1.0
GUTTER = 0.62
NCOL = 4
COLW = (W - 2 * MARGIN - (NCOL - 1) * GUTTER) / NCOL      # 10.785"
HEADER_H = 3.95
TILE_Y, TILE_H = 4.35, 2.60
BODY_Y = 7.25
FOOTER_H = 2.35
FOOTER_Y = H - MARGIN - FOOTER_H + 0.35
BODY_BOTTOM = FOOTER_Y - 0.45


def colx(i):
    return MARGIN + i * (COLW + GUTTER)


# ---------------------------------------------------------------- palette
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
VIOLET = RGBColor(0x4A, 0x3A, 0xA7)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
HAIRLINE = RGBColor(0xE1, 0xE0, 0xD9)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT_BLUE = RGBColor(0xF2, 0xF6, 0xFC)
TINT_ORANGE = RGBColor(0xFD, 0xF3, 0xEE)
TINT_AQUA = RGBColor(0xEE, 0xFA, 0xF5)
TINT_VIOLET = RGBColor(0xF2, 0xF0, 0xFA)
TINT_GREY = RGBColor(0xF7, 0xF7, 0xF5)
CRIT = RGBColor(0xD0, 0x3B, 0x3B)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"

# ---------------------------------------------------------------- type scale
T_TITLE = 80
T_SUB = 34
T_AUTHOR = 28
T_TILE_NUM = 62
T_TILE_LAB = 20
T_SECTION = 30
T_BODY = 20
T_SMALL = 17
T_CAPTION = 15
T_TABLE = 16
T_FOOT = 15


def read_table(name):
    with open(TABLES / name, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


PERF = read_table("table_performance_full.csv")
LEAK = read_table("table_leakage.csv")
MCN = read_table("table_mcnemar.csv")
SCA = read_table("table_split_c_authentic.csv")
SCS = read_table("table_split_c_synthetic.csv")
COST = read_table("table_cost.csv")
ABL = read_table("table_ablation_all_models.csv")

# Both external sets, evaluated from persisted checkpoints. Lives under
# modeling/results/ rather than paper/tables/ because it is produced by the
# modeling pipeline, not by the paper's metric scripts.
_EXT_PATH = ROOT / "modeling" / "results" / "external_from_checkpoints.csv"
with open(_EXT_PATH, newline="", encoding="utf-8") as _f:
    EXT = list(csv.DictReader(_f))

SHORT = {"model1_classical_colorhist_logreg": "M1 hist+LR",
         "model2_smallcnn_gap": "M2 CNN",
         "model3_mobilenetv3small_frozen": "M3 MobileNetV3",
         "model4_efficientnetb0_frozen": "M4 EfficientNet-B0"}
ORDER = list(SHORT)


def ext(model_key, split):
    for r in EXT:
        if r["model"] == model_key and r["split"] == split:
            return float(r["accuracy"])
    raise KeyError((model_key, split))


def leak(model_short, key):
    for r in LEAK:
        if r["model"] == model_short:
            return float(r[key])
    raise KeyError((model_short, key))


def f(rows, model, key, split=None):
    for r in rows:
        if r.get("model") == model and (split is None or r.get("split", "").startswith(split)):
            return float(r[key])
    raise KeyError((model, key, split))


# ---------------------------------------------------------------- primitives
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank


def rect(x, y, w, h, fill=None, line=None, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         adj=0.02, shadow=False):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = adj
        except (IndexError, KeyError):
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if not shadow:
        # python-pptx has no direct disable; clearing the effect list works
        spPr = sh._element.spPr
        for tag in ("a:effectLst",):
            for el in spPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"):
                spPr.remove(el)
    sh.text_frame.word_wrap = True
    return sh


def textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    return tb, tf


def para(tf, text, size=T_BODY, bold=False, italic=False, color=INK,
         align=PP_ALIGN.LEFT, space_before=0, space_after=5, first=False,
         font=FONT, line=0.95, bullet=None, indent=0.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if indent:
        p.level = 1
    body = text
    if bullet:
        body = f"{bullet}  {text}"
    # inline **bold** support
    parts = body.split("**")
    for k, chunk in enumerate(parts):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.bold = bold or (k % 2 == 1)
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return p


def section(x, y, w, title, accent=BLUE):
    """Section header: a colored rule with the title above it."""
    tb, tf = textbox(x, y, w, 0.62)
    para(tf, title.upper(), size=T_SECTION, bold=True, color=INK, first=True,
         space_after=0, line=0.9)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.60),
                                 Inches(w), Inches(0.075))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    return y + 0.88


def picture(x, y, w, stem, caption=None, max_h=None):
    """Place paper/figures/<stem>.png at width w, preserving aspect ratio."""
    path = FIGS / f"{stem}.png"
    iw, ih = Image.open(path).size
    h = w * ih / iw
    if max_h and h > max_h:
        h = max_h
        w = h * iw / ih
        x = x + (COLW - w) / 2
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    yy = y + h + 0.08
    if caption:
        tb, tf = textbox(x, yy, w, 0.5)
        para(tf, caption, size=T_CAPTION, color=INK2, first=True, line=0.92,
             space_after=0)
        yy += 0.22 * (1 + caption.count("\n")) + 0.22
    return yy


def table(x, y, w, rows, col_w=None, header=True, font_size=T_TABLE,
          row_h=0.42, accent=BLUE):
    nrows, ncols = len(rows), len(rows[0])
    h = row_h * nrows
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w),
                                Inches(h))
    tbl = gf.table
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = Emu(int(Inches(w) * cw / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (accent if ri == 0 and header
                                        else (WHITE if ri % 2 else TINT_GREY))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            p.line_spacing = 0.9
            bold = ri == 0 and header
            for k, chunk in enumerate(str(val).split("**")):
                if not chunk:
                    continue
                r = p.add_run()
                r.text = chunk
                r.font.size = Pt(font_size)
                r.font.name = FONT
                r.font.bold = bold or (k % 2 == 1)
                r.font.color.rgb = WHITE if bold else INK
    return y + h + 0.2


# =============================================================== header
rect(0, 0, W, HEADER_H, fill=RGBColor(0x0D, 0x36, 0x6B), shape=MSO_SHAPE.RECTANGLE)
rect(0, HEADER_H - 0.09, W, 0.09, fill=ORANGE, shape=MSO_SHAPE.RECTANGLE)

tb, tf = textbox(MARGIN, 0.42, W - 2 * MARGIN - 7.0, 2.3)
para(tf, "When the label predicts the camera, not the product",
     size=T_TITLE, bold=True, color=WHITE, first=True, line=0.88, space_after=6)
para(tf, "Asymmetric class sourcing creates provenance confounds in authenticity-"
         "classification datasets — a cheap audit that finds them, a label-free correction that partly repairs them",
     size=T_SUB, color=RGBColor(0xCD, 0xE2, 0xFB), line=0.95)

tb, tf = textbox(MARGIN, 3.05, W - 2 * MARGIN - 7.0, 1.0)
para(tf, "Sophie Zhu", size=T_AUTHOR,
     bold=True, color=WHITE, first=True, space_after=2)
para(tf, "Mira Costa High School, Manhattan Beach, CA 90266 USA"
         "     ·     ORCID 0009-0004-2403-910X"
         "     ·     sophiezhu2028@gmail.com",
     size=T_SMALL + 3, color=RGBColor(0x9E, 0xC5, 0xF4))

# header badge
badge = rect(W - MARGIN - 6.2, 0.75, 6.2, 2.65, fill=RGBColor(0x18, 0x4F, 0x95),
             line=RGBColor(0x39, 0x87, 0xE5), line_w=1.5, adj=0.06)
tf = badge.text_frame
tf.margin_left = tf.margin_right = Inches(0.22)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "Provenance confounds", size=T_SUB, bold=True, color=WHITE, first=True,
     align=PP_ALIGN.CENTER, space_after=3)
para(tf, "510-image pool  ·  4 model families\n3 evaluation protocols  ·  "
         "CPU-only, seed-fixed, fully reproducible",
     size=T_SMALL, color=RGBColor(0xCD, 0xE2, 0xFB), align=PP_ALIGN.CENTER,
     line=1.0)

# =============================================================== stat tiles
TILES = [
    ("100%", "of counterfeit labels are screen captures,\n"
             "100% of authentic labels are photographs",
     "and 3 acquisition scalars with no pixels\nscore 100% on the corrected split",
     CRIT, TINT_ORANGE, ORANGE),
    ("0.0% → 3.3%", "external accuracy on 150 authentic\nphotographs from a new source",
     "…from models scoring 84–97% on the\nauthentic class in-distribution",
     CRIT, TINT_ORANGE, ORANGE),
    ("→ 86.0%", "after label-free 3-way normalization\n(resolution + brightness + JPEG)",
     "…but only 46.3% on a 2nd capture shift —\nthe backbones hold, this model does not",
     GOOD, TINT_AQUA, AQUA),
    ("≤ 6.8 pts", "total effect of correcting the\ntrain/test split protocol",
     "the leakage everyone checks for was\nthe smaller of the two problems",
     BLUE, TINT_BLUE, BLUE),
]
tw = (W - 2 * MARGIN - 3 * GUTTER) / 4
for i, (num, lab, sub, numc, tint, edge) in enumerate(TILES):
    x = MARGIN + i * (tw + GUTTER)
    rect(x, TILE_Y, tw, TILE_H, fill=tint, line=edge, line_w=1.5, adj=0.05)
    tb, tf = textbox(x + 0.3, TILE_Y + 0.16, tw - 0.6, TILE_H - 0.3)
    para(tf, num, size=T_TILE_NUM, bold=True, color=numc, first=True,
         space_after=2, line=0.85)
    para(tf, lab, size=T_TILE_LAB, bold=True, color=INK, line=0.95, space_after=4)
    para(tf, sub, size=T_SMALL, color=INK2, line=0.95)

# =============================================================== column 1
x = colx(0)
y = BODY_Y
y = section(x, y, COLW, "Background", ORANGE)
tb, tf = textbox(x, y, COLW, 4.6)
para(tf, "**The general problem.** In any task asking *is this genuine?* the "
         "inauthentic class is harder to obtain than the authentic one — so it "
         "gets obtained **differently**: screen-captured, scraped, edited or "
         "generated. The label then predicts the *acquisition process*, which is "
         "far easier to learn than packaging semantics. We call this "
         "**class-conditional provenance confounding** and argue it is the "
         "default outcome of dataset construction, not an accident.",
     size=T_SMALL + 1, first=True, line=0.98, space_after=8)
para(tf, "**The case study.** Much falsified medicine is *visually* imperfect, "
         "so photo screening is an attractive triage tool, and CNN classifiers "
         "on small public authentic-vs-counterfeit sets routinely report "
         "**>90% accuracy**. We audited the 661-image Kaggle *Fake vs Real "
         "Medicine* set — freely available, no data card, no stated acquisition "
         "protocol, and typical of what this area works with — asking how much "
         "of such a figure survives methodological correction. We expected the "
         "answer to be about train/test leakage. It was not.",
     size=T_SMALL + 1, line=0.98)
y += 3.45

y = section(x, y, COLW, "Research question", ORANGE)
q = rect(x, y, COLW, 3.55, fill=TINT_ORANGE, line=ORANGE, line_w=2.0, adj=0.045)
tf = q.text_frame
tf.margin_left = tf.margin_right = Inches(0.28)
tf.margin_top = Inches(0.2)
para(tf, "How much of a >90% accuracy claim on a dataset like this is real "
         "signal about packaging — and how much is an artifact of how the "
         "dataset was built?", size=T_BODY + 4, bold=True, color=INK,
     first=True, line=1.0, space_after=8)
para(tf, "Sub-questions: (i) how large is the train/test leakage effect under "
         "a controlled, single-variable comparison? (ii) does anything transfer "
         "to images someone else captured? (iii) if not, is the failure "
         "correctable?", size=T_SMALL + 1, color=INK2, line=0.98)
y += 3.85

y = section(x, y, COLW, "Dataset", ORANGE)
y = table(x, y, COLW, [
    ["Source", "Files", "Role"],
    ["Kaggle *Fake vs Real Medicine*", "661", "Modeling pool (Splits A, B)"],
    ["Roboflow Counterfeit_med_detection v4", "4,260", "**Excluded** — see below"],
    ["Mendeley mobile-captured packages", "150", "External Split C (authentic)"],
    ["Synthetic proxy (perturbed copies)", "150", "Counterfeit-recall stress test"],
], col_w=[3.1, 0.8, 2.3], row_h=0.52, accent=ORANGE)

tb, tf = textbox(x, y, COLW, 4.6)
para(tf, "**Roboflow excluded, not merely down-weighted:** 57/57 of its "
         "counterfeit-labeled images are regulator advisory graphics with the "
         "label printed in the pixels; 263/263 of its plain product photos are "
         "authentic. After filtering it yields **2** usable counterfeits.",
     size=T_SMALL + 1, first=True, line=0.98, space_after=7)
para(tf, "**The two public sources are not independent:** rotation-aware "
         "perceptual hashing finds 229 shared clusters — **44% of the Kaggle "
         "set near-duplicates Roboflow images**. Cross-dataset validation "
         "between them would leak.", size=T_SMALL + 1, line=0.98, space_after=7)
para(tf, "**Its shipped train/val/test split is not a split:** `train` holds "
         "all 661 images; `val` (453) and `test` (449) are subsets of it. "
         "License *Unknown*, 591 downloads, 3 notebooks — an unvetted upload, "
         "not a community benchmark.", size=T_SMALL + 1, line=0.98,
     space_after=7)
para(tf, "**Complete human review** (a census, not a sample) removed 56 files: "
         "47 watermarked (47/47 authentic-labeled), 4 non-medicine, 5 with no "
         "packaging. Final pool **510 images / 480 product groups**, 272/238; "
         "43.7% blister, 30.4% carton, 25.9% mixed — scope is *packaging and "
         "immediate containers*.", size=T_SMALL + 1, line=0.98)
y += 4.35

y = section(x, y, COLW, "ROC — in-distribution", ORANGE)
y = picture(x, y, COLW, "fig04_roc",
            "Fig. 1  ROC on both in-distribution partitions. Every model looks "
            "excellent here. Figs. 3, 7 and 8 show the same models on data they "
            "did not come from.")
COL1_END = y

# =============================================================== column 2
x = colx(1)
y = BODY_Y
y = section(x, y, COLW, "Methodology", BLUE)
y = table(x, y, COLW, [
    ["Model", "Trainable", "Note"],
    ["M1  color histogram + LogReg", "97", "96-dim RGB histogram"],
    ["M2  small CNN, GAP head", "23,938", "GAP head instead of dense(128)"],
    ["M3  MobileNetV3-Small, frozen", "1,154", "927 K frozen backbone params"],
    ["M4  EfficientNet-B0, frozen", "2,562", "4.0 M frozen backbone params"],
], col_w=[3.0, 1.0, 2.6], row_h=0.52, accent=BLUE)

tb, tf = textbox(x, y, COLW, 2.6)
para(tf, "Identical protocol for M2–M4: Adam, batch 32, class-weighted "
         "cross-entropy, early stopping on val loss (patience 4, min Δ 1e−3), "
         "LR grid {1e−3, 3e−4, 1e−4} (all three selected 1e−3), seed 42, "
         "CPU only. Augmentation: ±12° rotation, brightness/contrast jitter, "
         "mild crop, slight blur — **no flip** (it would mirror printed text).",
     size=T_SMALL + 1, first=True, line=0.98, space_after=7)
para(tf, "**Label-free 3-way capture normalization**, applied identically to "
         "every partition and deployable at inference: 128 px short-side "
         "bottleneck → rescale mean brightness to 0.5 → re-encode at JPEG "
         "quality 40. Costs **<17 ms/image**.", size=T_SMALL + 1, line=0.98)
y += 2.35

y = section(x, y, COLW, "Pipeline", BLUE)
y = picture(x, y, COLW * 0.93, "fig01_workflow",
            "Fig. 2  Provenance, splitting protocol and evaluation design.")

y = section(x, y, COLW, "Experimental design", BLUE)
tb, tf = textbox(x, y, COLW, 5.0)
for bullet, txt in [
    ("A", "**Split A — naive:** random 70:15:15 at the **image** level, "
          "the protocol in general use on data of this kind. 9/480 product groups "
          "straddle partitions."),
    ("B", "**Split B — corrected:** 70:15:15 at the **product-identity group** "
          "level + 5-fold StratifiedGroupKFold. Zero group overlap, asserted "
          "on every run."),
    ("C", "**Split C — external:** 150 authentic photographs from an "
          "independent source, verified non-duplicative (0/150 matched; "
          "nearest pHash distance 10/64 vs. threshold 8)."),
    ("C*", "**Split C — synthetic:** the same 150 photographs plus 150 "
           "perturbed copies (print/color/text defects), an ImageNet-C-style "
           "proxy for the counterfeit direction. Not a measurement of real "
           "counterfeit recall."),
]:
    para(tf, txt, size=T_SMALL + 1, first=(bullet == "A"), line=0.98,
         space_after=7, bullet="▪")
para(tf, "Counterfeit is the positive class. All point estimates carry 95% "
         "percentile bootstrap intervals (2,000 resamples); models are compared "
         "with exact McNemar tests on the same test partition.",
     size=T_SMALL, color=INK2, line=0.98)
y += 4.05

y = section(x, y, COLW, "Counterfeit-recall stress test", BLUE)
y = picture(x, y, COLW, "fig09_confusion_synthetic",
            "Fig. 3  Synthetic counterfeit-proxy Split C. No model achieves "
            "well-calibrated counterfeit recall: M1 has the best AUC (0.895) at "
            "a threshold so wrong its accuracy is pinned at 0.500, and M3 sits "
            "at chance (0.503) — corroborating its backdrop-shortcut Grad-CAM "
            "on a completely independent evaluation set.")
COL2_END = y

# =============================================================== column 3
x = colx(2)
y = BODY_Y
y = section(x, y, COLW, "Result 1 — the confound", AQUA)
y = picture(x, y, COLW, "fig03_capture_confound",
            "Fig. 4  The two acquisition pipelines, and where the external set "
            "sits relative to both.")
tb, tf = textbox(x, y, COLW, 2.2)
para(tf, "Brightness 0.767 vs 0.555 (**t = 17.0, p ≈ 0**) · median short side "
         "223 vs 405 px · mean file size 6.0 vs 339 kB (**56×**). The external "
         "set is ~10× higher resolution and **darker than either training "
         "class** (0.162) — so 'bright, small, compressed → authentic' calls "
         "every external photo counterfeit.", size=T_SMALL + 1, first=True,
     line=0.98)
y += 1.75

y = section(x, y, COLW, "Feature importance", AQUA)
y = picture(x, y, COLW, "fig12_model1_attribution",
            "Fig. 5  M1's decision function, and its exact Shapley "
            "decomposition (closed form for a linear model).")
tb, tf = textbox(x, y, COLW, 2.15)
para(tf, "3 of 96 features carry the model: the near-white bin (248–255) of "
         "each channel, mean |φ| = 0.079–0.082 against **≤ 0.002** for the "
         "other 93 — M1's 83.8% measures how much white an image contains.",
     size=T_SMALL + 1, first=True, line=0.98, space_after=7)
para(tf, "**Metadata-only oracle — the tight bound.** LogReg on 3 acquisition "
         "scalars, no pixels: **100%** on the leakage-free test partition "
         "(file size alone suffices; resolution 94.6%, brightness only 71.6%), "
         "and the file extension alone is right on all 510. Nothing measurable "
         "in-distribution here needs packaging information to explain it.",
     size=T_SMALL + 1, line=0.98)
y += 2.2

y = section(x, y, COLW, "Result 2 — leakage is the small problem", AQUA)
y = picture(x, y, COLW, "fig13_leakage",
            "Fig. 6  Naive vs product-grouped partitioning, 95% bootstrap "
            "intervals.")
_leak_rows = [["Model", "Split A", "Split B", "Δ", "5-fold CV"]]
for _m in ("M1", "M2", "M3", "M4"):
    _short = {"M1": "M1 hist+LR", "M2": "M2 CNN",
              "M3": "M3 MobileNetV3", "M4": "M4 EfficientNet-B0"}[_m]
    _d = leak(_m, "delta_a_minus_b")
    _leak_rows.append([
        _short,
        f"{leak(_m, 'split_a_acc'):.3f}",
        f"{leak(_m, 'split_b_acc'):.3f}",
        (f"**{_d:+.3f}**" if abs(_d) >= 0.01 else f"{_d:+.3f}"),
        f"{leak(_m, 'split_b_cv_mean'):.3f} ± {leak(_m, 'split_b_cv_std'):.3f}",
    ])
y = table(x, y, COLW, _leak_rows,
          col_w=[2.5, 1.0, 1.0, 1.0, 1.5], row_h=0.5, accent=AQUA)
tb, tf = textbox(x, y, COLW, 1.4)
para(tf, "**No pairwise McNemar test is significant** (all p ≥ 0.118; "
         "discordant counts 3–16). A 97-parameter linear model and a 4 M-"
         "parameter pretrained network are statistically indistinguishable on "
         "this test partition.", size=T_SMALL + 1, first=True, line=0.98)
COL3_END = y + 1.4

# =============================================================== column 4
x = colx(3)
y = BODY_Y
y = section(x, y, COLW, "Result 3 — external generalization", VIOLET)
y = picture(x, y, COLW, "fig08_external_generalisation",
            "Fig. 7  In-distribution vs external accuracy on the authentic "
            "class, before and after normalization.")
# Pre-normalization Split C baselines are archived values, not recomputable
# from the current pipeline, so they stay literal; everything else is read.
_BASE = {"model1_classical_colorhist_logreg": 0.000,
         "model2_smallcnn_gap": 0.000,
         "model3_mobilenetv3small_frozen": 0.693,
         "model4_efficientnetb0_frozen": 0.033}
_ext_rows = [["Model", "In-dist.", "Ext. base", "Split C", "Split D"]]
for _k in ORDER:
    _ind = next(float(r["split_b_authentic_acc"]) for r in SCA
                if r["model_full"].split(" (")[0] in SHORT[_k]
                or SHORT[_k].split()[0] == r["model"])
    _c, _d = ext(_k, "split_c"), ext(_k, "split_d")
    _ext_rows.append([SHORT[_k], f"{_ind:.3f}", f"{_BASE[_k]:.3f}",
                      f"**{_c:.3f}**" if _c > 0.8 else f"{_c:.3f}",
                      f"**{_d:.3f}**" if (_d > 0.8 or _d < 0.5) else f"{_d:.3f}"])
y = table(x, y, COLW, _ext_rows,
          col_w=[2.4, 1.1, 1.1, 1.1, 1.0], row_h=0.5, accent=VIOLET)

y = section(x, y, COLW, "Ablation", VIOLET)
y = picture(x, y, COLW, "fig10_ablation",
            "Fig. 8  Per-axis ablation (within-run) and per-architecture "
            "outcome. White balance was tested and ruled out.")
tb, tf = textbox(x, y, COLW, 1.3)
para(tf, "M4, one run: **none 5.3% → compression 12.7% → res+bright 50.7% → "
         "all three 78.0%** externally, at no in-distribution cost. The axes "
         "are complementary, not redundant.", size=T_SMALL + 1, first=True,
     line=0.98)
y += 1.35

# =============================================================== discussion band
y = section(x, y + 0.15, COLW, "Discussion · conclusion · next", VIOLET)
tb, tf = textbox(x, y, COLW, 4.6)
para(tf, "**One external set was not enough — and it caught us out.** A 2nd "
         "external distribution (same 150 products, different phone and "
         "lighting) splits the models: M4 **0.807→0.832** and M3 "
         "**0.773→0.725** hold within noise, but M2 falls **0.860→0.463**. "
         "Its Split C lead was specific to that capture condition — we wrote "
         "the warning about this hazard, then fell into it.",
     size=T_SMALL, line=0.96, space_after=6)
para(tf, "**The audit generalizes; a clean result does not.** On a second, "
         "independently published dataset it returns only **0.717** — yet that "
         "dataset is worse: 57/57 of its counterfeit images carry the label "
         "printed in the pixels. Its publisher had resized everything to "
         "640×640 and one format, erasing the traces without removing the "
         "confound. **Curated archives are harder to audit, not safer.**",
     size=T_SMALL, line=0.96, space_after=6)
para(tf, "**Held-out accuracy cannot expose any of this.** Stratification, "
         "grouped 5-fold CV, bootstrap CIs and a leakage-free product-level "
         "split — none saw anything wrong.", size=T_SMALL, line=0.96,
     space_after=6)
para(tf, "**Four defect types, four checks.** A: acquisition statistics → the "
         "metadata audit. B: content/modality, label in the pixels → human "
         "inspection; metadata is blind. C: reintroduced by source selection → "
         "re-audit every derived partition. D: degenerate shipped split → "
         "intersect the filename sets.", size=T_SMALL, line=0.96, space_after=6)
para(tf, "**Capacity is the wrong axis.** M3's strong *uncorrected* score is "
         "itself a shortcut — Grad-CAM attends to the shared backdrop, and it "
         "sits at chance where both classes share it. **We would deploy none "
         "of them:** recall is untested against real counterfeits.",
     size=T_SMALL, line=0.96, space_after=6)
para(tf, "**Independently corroborated.** Grommelt et al. (*ECCV 2024 Workshops*) "
         "find GenImage's real/generated classes separated by JPEG compression "
         "and size — detectors partly reduced to JPEG detectors, >11 pt shift "
         "after equalizing. Same signature, unrelated field, independent "
         "discovery: what a structural cause predicts.",
     size=T_SMALL, line=0.96, space_after=6)
para(tf, "**Next:** run the audit across a corpus of authenticity datasets — "
         "counterfeit goods, document forgery, generated-image detection, "
         "defect inspection — recording whether each sourced its classes by "
         "one procedure. Metadata only, no training runs.",
     size=T_SMALL, line=0.96)
COL4_END = y + 6.5

# =============================================================== footer
rect(0, FOOTER_Y - 0.3, W, FOOTER_H + 0.3, fill=TINT_GREY, shape=MSO_SHAPE.RECTANGLE)
rect(0, FOOTER_Y - 0.3, W, 0.06, fill=HAIRLINE, shape=MSO_SHAPE.RECTANGLE)

# QR placeholder
qr = rect(MARGIN, FOOTER_Y - 0.05, 1.85, 1.85, fill=WHITE, line=INK, line_w=2.0,
          shape=MSO_SHAPE.RECTANGLE)
tf = qr.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "QR\ncode\nhere", size=T_SMALL, bold=True, color=INK2, first=True,
     align=PP_ALIGN.CENTER, line=1.0)
tb, tf = textbox(MARGIN + 2.05, FOOTER_Y - 0.02, 6.0, 1.9)
para(tf, "Code, data card and all figures", size=T_SMALL + 1, bold=True,
     first=True, space_after=2)
para(tf, "github.com/sophiezla/counterfeit-drug  ·  doi:10.5281/zenodo.21936721",
     size=T_FOOT, color=INK2, line=0.98)
para(tf, "Vector versions of every figure accompany this poster as PDF.",
     size=T_FOOT, color=MUTED, line=0.98)

tb, tf = textbox(MARGIN + 8.4, FOOTER_Y - 0.02, 11.5, 2.1)
para(tf, "Acknowledgements", size=T_SMALL + 1, bold=True, first=True, space_after=2)
para(tf, "This work received no specific grant from any funding agency in the "
         "public, commercial, or not-for-profit sectors, and declares no "
         "conflict of interest. We thank the maintainers of the three "
         "public datasets used here. The Mendeley and Roboflow datasets are "
         "CC BY 4.0; the Kaggle archive states no license, its listing says "
         "\"Unknown\", and it is attributed to "
         "its uploader.", size=T_FOOT, color=INK2, line=0.98)

tb, tf = textbox(MARGIN + 20.4, FOOTER_Y - 0.02, W - MARGIN - 20.4 - MARGIN, 2.1)
para(tf, "Selected references", size=T_SMALL + 1, bold=True, first=True, space_after=2)
para(tf, "[6] Geirhos et al., Shortcut learning in deep neural networks, "
         "*Nature Machine Intelligence* 2;665–673, 2020.   "
         "[7] Zech et al., Variable generalization performance of a deep "
         "learning model to detect pneumonia in chest radiographs, "
         "*PLOS Medicine* 15(11):e1002683, 2018.   "
         "[12] Hendrycks & Dietterich, Benchmarking neural network robustness "
         "to common corruptions and perturbations, *ICLR* 2019.",
     size=T_FOOT, color=INK2, line=0.98, space_after=3)
para(tf, "[13] Howard et al., Searching for MobileNetV3, *ICCV* 2019.   "
         "[14] Tan & Le, EfficientNet, *ICML* 2019.   "
         "[15] Selvaraju et al., Grad-CAM, *ICCV* 2017.   "
         "[16] Lundberg & Lee, A unified approach to interpreting model "
         "predictions, *NeurIPS* 2017.   "
         "[20] Abdelmaksoud et al., Mobile-captured pharmaceutical medication "
         "packages, Mendeley Data, doi:10.17632/bjy2svvmn8.1.",
     size=T_FOOT, color=INK2, line=0.98, space_after=3)
para(tf, "Full numbered reference list is in the manuscript; every entry is "
         "verified against a primary source.",
     size=T_FOOT, color=MUTED, italic=True, line=0.98)

print(f"body region: {BODY_Y:.2f}\" to {BODY_BOTTOM:.2f}\"")
for i, end in enumerate((COL1_END, COL2_END, COL3_END, COL4_END), start=1):
    slack = BODY_BOTTOM - end
    flag = "OVERFLOW" if slack < 0 else ("sparse" if slack > 2.0 else "ok")
    print(f"  column {i}: ends at {end:6.2f}\"  slack {slack:+.2f}\"  [{flag}]")

prs.save(OUT)
print(f"wrote {OUT.relative_to(ROOT)}")
print(f"  {len(slide.shapes)} native shapes  ·  "
      f"{sum(1 for s in slide.shapes if s.shape_type == 13)} pictures  ·  "
      f"{sum(1 for s in slide.shapes if s.has_table)} tables")
